from pptx.util import Inches
from pptx import Presentation

# Redefine the presentation object to avoid duplication
prs = Presentation()

# Recreate title and intro slides
add_title_slide("Data Analytics Case Study – College Placement Analysis", "Your Name | Class | Date")
add_content_slide("Objective of the Analysis", [
    "Understand placement trends in the college.",
    "Analyze factors that influence student placement.",
    "Provide insights to improve placement strategies."
])
add_content_slide("About the Dataset", [
    "Student information: Gender, Stream, Internships, CGPA, Hostel, Backlogs, Placement.",
    "Total entries: 2,966 students.",
    "Target variable: PlacedOrNot (1 = Placed, 0 = Not Placed)."
])

# Function to add image and analysis side by side (without using MSO_ANCHOR)
def add_image_with_simple_analysis(title, image_path, analysis_points):
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Image
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), width=Inches(4.5))

    # Textbox for analysis
    textbox = slide.shapes.add_textbox(Inches(5.2), Inches(1), Inches(4), Inches(4))
    tf = textbox.text_frame
    for i, point in enumerate(analysis_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {point}"
        p.font.size = Pt(14)

# Reuse the same analysis as earlier
for title, img, points in image_analyses:
    add_image_with_simple_analysis(title, img, points)

# Final summary slides
add_content_slide("Insights & Observations", [
    "Internships greatly improve placement chances.",
    "High CGPA correlates with placement.",
    "Backlogs and lack of internships reduce chances.",
    "Stream plays a role in opportunities."
])

add_content_slide("Conclusion", [
    "Academic performance and internships are key drivers for placements.",
    "Career readiness can be improved through practical exposure.",
    "Insights help departments in guiding students better."
])

add_content_slide("Any Questions?", ["Feel free to ask!"])

# Save updated presentation
pptx_fixed = "/mnt/data/College_Placement_Analysis_With_Insights_Fixed.pptx"
prs.save(pptx_fixed)

pptx_fixed
